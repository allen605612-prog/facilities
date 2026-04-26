#!/usr/bin/env python3
"""
第 66 屆科展得獎橫幅
Banner ratio from PPTX: cx=40684450, cy=12780963 EMU → ≈ 3.183:1
"""
from PIL import Image, ImageDraw, ImageFont

W, H = 5100, 1602

BG_DARK  = (11,  38,  24)
BG_MID   = (16,  52,  33)
BG_CARD  = (14,  46,  28)
GOLD     = (212, 175,  55)
GOLD_LT  = (255, 223,  96)
GOLD_DIM = (140, 108,  22)
CREAM    = (255, 250, 215)
GRN      = (22,  78,  44)

WDIR = r"C:\Windows\Fonts"
CDIR = r"C:\Users\user\.claude\plugins\cache\anthropic-agent-skills\document-skills\b0cbd3df1533\skills\canvas-design\canvas-fonts"

def ttf(path, size):
    try:    return ImageFont.truetype(path, size)
    except: return ImageFont.load_default()

f_header = ttf(rf"{WDIR}\msjhbd.ttf",  84)
f_huge   = ttf(rf"{WDIR}\msjhbd.ttf", 190)
f_sec    = ttf(rf"{WDIR}\msjhbd.ttf",  60)
f_award  = ttf(rf"{WDIR}\msjhbd.ttf",  54)
f_card   = ttf(rf"{WDIR}\msjhbd.ttf",  44)
f_body   = ttf(rf"{WDIR}\msjhbd.ttf",  40)
f_sm     = ttf(rf"{WDIR}\msjh.ttf",    34)

img  = Image.new("RGBA", (W, H), (0, 0, 0, 255))
draw = ImageDraw.Draw(img)

# Background + glow
draw.rectangle([0, 0, W, H], fill=BG_DARK)
for i in range(40):
    v = int(BG_MID[1] * (1 - i/40))
    draw.rectangle([80+i, 80+i, W-80-i, H-80-i], outline=(BG_MID[0], v, BG_MID[2]), width=1)

# Triple border
for m, w, c in [(20, 5, GOLD), (32, 1, GOLD_DIM), (46, 3, GOLD)]:
    draw.rectangle([m, m, W-m, H-m], outline=c, width=w)

def diamond(cx, cy, r=18):
    draw.polygon([(cx,cy-r),(cx+r*3//5,cy),(cx,cy+r),(cx-r*3//5,cy)], fill=GOLD)
    draw.polygon([(cx,cy-r//2),(cx+r//4,cy),(cx,cy+r//2),(cx-r//4,cy)], fill=BG_DARK)

for cx_, cy_ in [(56, 56), (W-56, 56), (56, H-56), (W-56, H-56)]:
    diamond(cx_, cy_)

# ── HEADER BAR ────────────────────────────────────────────────────────
HDR_Y0, HDR_H = 56, 148
draw.rectangle([56, HDR_Y0, W-56, HDR_Y0+HDR_H], fill=GRN)
draw.line([(56, HDR_Y0+HDR_H), (W-56, HDR_Y0+HDR_H)], fill=GOLD, width=3)

hdr = "雲林縣私立正心高級中學  ·  科展榮耀榜"
bb  = draw.textbbox((0,0), hdr, font=f_header)
draw.text(((W-(bb[2]-bb[0]))//2, HDR_Y0+(HDR_H-(bb[3]-bb[1]))//2-bb[1]), hdr,
          font=f_header, fill=GOLD_LT)

# ── LAYOUT ────────────────────────────────────────────────────────────
CONT_Y = HDR_Y0 + HDR_H + 16
CONT_B = H - 56
AVAIL  = CONT_B - CONT_Y    # ≈ 1322

DIV_X  = 1660

draw.line([(DIV_X,   CONT_Y+10), (DIV_X,   CONT_B-10)], fill=GOLD,    width=2)
draw.line([(DIV_X+7, CONT_Y+20), (DIV_X+7, CONT_B-20)], fill=GOLD_DIM, width=1)

# ══════════════════════════════════════════════════════════════════════
# LEFT — 高中科展  (vertically centred)
# ══════════════════════════════════════════════════════════════════════
def text_h(text, font):
    bb = draw.textbbox((0,0), text, font=font)
    return bb[3]-bb[1]

# Measure left block height
huge_h  = text_h("狂賀！", f_huge)
sub_h   = text_h("第 66 屆第四區分區高中科展", f_sec)
at_h    = text_h("榮獲「環境學科」優等", f_award)
band_h  = at_h + 26
stud_h  = text_h("高一心班 宋彥霖", f_body)
teach_h = text_h("指導老師：張祐誠老師", f_body)

LEFT_BLOCK_H = (huge_h + 14 + sub_h + 16 + band_h + 18 + stud_h + 12 + teach_h)
LEFT_Y0  = CONT_Y + (AVAIL - LEFT_BLOCK_H) // 2
LEFT_CX  = (56 + DIV_X) // 2

y = LEFT_Y0
# 狂賀！
bb_h = draw.textbbox((0,0), "狂賀！", font=f_huge)
hx = LEFT_CX - (bb_h[2]-bb_h[0])//2
draw.text((hx+5, y+5), "狂賀！", font=f_huge, fill=GOLD_DIM)
draw.text((hx,   y),   "狂賀！", font=f_huge, fill=GOLD_LT)
y += huge_h + 14

# Sub title
draw.text((LEFT_CX, y), "第 66 屆第四區分區高中科展", font=f_sec, fill=CREAM, anchor='mt')
y += sub_h + 16

# Award band
at_text = "榮獲「環境學科」優等"
bb_at   = draw.textbbox((0,0), at_text, font=f_award)
aw = bb_at[2]-bb_at[0]
bx0, bx1 = LEFT_CX-aw//2-36, LEFT_CX+aw//2+36
draw.rectangle([bx0, y, bx1, y+band_h], fill=GOLD, outline=GOLD_DIM, width=2)
draw.text((LEFT_CX, y+band_h//2), at_text, font=f_award, fill=BG_DARK, anchor='mm')
y += band_h + 18

# Students
draw.text((LEFT_CX, y), "高一心班 宋彥霖・高一意班 鐘宥昕・高一正班 廖柃柃",
          font=f_body, fill=CREAM, anchor='mt')
y += stud_h + 12

# Teacher
draw.text((LEFT_CX, y), "指導老師：張祐誠老師", font=f_body, fill=GOLD, anchor='mt')

# ══════════════════════════════════════════════════════════════════════
# RIGHT — 縣科展  (cards with fixed height, vertically centred block)
# ══════════════════════════════════════════════════════════════════════
RX0 = DIV_X + 22
RW  = W - RX0 - 56
RCX = RX0 + RW // 2

# Section header
sec = "雲林縣第 66 屆公私立國民中小學科學展覽會"
bb_sec = draw.textbbox((0,0), sec, font=f_sec)
sec_h_px = bb_sec[3]-bb_sec[1]

# Card dimensions — fixed height so content sits tight
CARD_H   = 210
CARD_GAP = 14
CARDS_BLOCK = 2*CARD_H + CARD_GAP
TOTAL_BLOCK = sec_h_px + 22 + CARDS_BLOCK
BLOCK_Y0    = CONT_Y + (AVAIL - TOTAL_BLOCK) // 2

draw.text((RCX, BLOCK_Y0), sec, font=f_sec, fill=GOLD_LT, anchor='mt')
CARDS_Y0 = BLOCK_Y0 + sec_h_px + 22

# Card width for 4-card row
CARD_GAP3 = 14
card_w = (RW - 3*CARD_GAP3) // 4

RANK_CLR = {
    "第一名": (185, 150, 18),
    "第二名": (158, 126, 13),
    "第三名": (128, 100, 12),
    "第四名": (105,  80, 12),
}

awards = [
    ("數學科",     "第一名", "廖采玲  陳彙喬  黃柏之", "羅婉萍老師"),
    ("數學科",     "第二名", "廖泱泱",              "詩富強老師"),
    ("生物科",     "第二名", "洪靖宇  蔡承恩",      "許議鶴老師"),
    ("生物科",     "第三名", "黃家恩  王嘉裕  陳昱勳","郭靜如老師"),
    ("生活應用(一)","第四名", "陳沛青  鍾子鵬  陳意凡","林俊杰老師"),
    ("生活應用(二)","第三名", "籃苡熏  萬奕岑  陳姸汝","許議鶴老師"),
    ("生活應用(三)","第四名", "許呈睿  黃品語  洪啟鈞","許議鶴老師"),
]

for i, (subj, rank, students, teacher) in enumerate(awards):
    if i < 4:
        cx0 = RX0 + i*(card_w + CARD_GAP3)
        cy0 = CARDS_Y0
    else:
        col = i - 4
        row2_total = 3*card_w + 2*CARD_GAP3
        row2_x0 = RX0 + (RW - row2_total)//2
        cx0 = row2_x0 + col*(card_w + CARD_GAP3)
        cy0 = CARDS_Y0 + CARD_H + CARD_GAP

    cx1 = cx0 + card_w
    cy1 = cy0 + CARD_H
    ccx = (cx0+cx1)//2

    # Card body
    draw.rectangle([cx0, cy0, cx1, cy1], fill=BG_CARD, outline=GOLD_DIM, width=1)

    # Award band
    BAND_H = 62
    draw.rectangle([cx0, cy0, cx1, cy0+BAND_H], fill=RANK_CLR.get(rank, GOLD_DIM))
    draw.text((ccx, cy0+BAND_H//2), f"{subj}  {rank}", font=f_card, fill=GOLD_LT, anchor='mm')

    # Students
    s_y = cy0 + BAND_H + 14
    draw.text((ccx, s_y), students, font=f_sm, fill=CREAM, anchor='mt')
    bb_s = draw.textbbox((0,0), students, font=f_sm)
    line_y = s_y + (bb_s[3]-bb_s[1]) + 12
    draw.line([(cx0+24, line_y), (cx1-24, line_y)], fill=GOLD_DIM, width=1)

    # Teacher (anchored near bottom of card, not at very edge)
    draw.text((ccx, cy1-16), teacher, font=f_sm, fill=GOLD, anchor='mb')

# ── Save ──────────────────────────────────────────────────────────────
out_img  = img.convert("RGB")
out_png  = r"C:\Users\user\allen\第66屆科展得獎橫幅.png"
out_pptx = r"C:\Users\user\allen\第66屆科展得獎橫幅.pptx"
out_img.save(out_png, "PNG")

from pptx import Presentation
from pptx.util import Emu
prs = Presentation()
prs.slide_width  = Emu(40684450)
prs.slide_height = Emu(12780963)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(out_png, Emu(0), Emu(0), Emu(40684450), Emu(12780963))
prs.save(out_pptx)

print(f"PNG  → {out_png}")
print(f"PPTX → {out_pptx}")
print("Done.")
