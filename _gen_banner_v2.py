#!/usr/bin/env python3
"""
第 66 屆科展得獎橫幅 v2
  活潑版: 白底彩旗活潑色調
  照片版: 嵌入各得獎照片
"""
import math
from PIL import Image, ImageDraw, ImageFont

W, H = 5100, 1602

# ── Palette — 活潑 ─────────────────────────────────────────────────
BG       = (252, 252, 248)
NAVY     = (22,  55, 130)
DARK_GRN = (18,  88,  48)
TEAL     = (0,  152, 138)
GREEN    = (38, 158,  72)
CORAL    = (218,  72,  52)
GOLD     = (205, 160,  18)
ORANGE   = (222, 128,  28)
WHITE    = (255, 255, 255)
DARK     = (32,  32,  32)
CREAM    = (255, 252, 228)

BUNTING_COLORS = [CORAL, GOLD, GREEN, NAVY, (132, 58, 188), TEAL, ORANGE]

SUBJ_COLORS = {
    "數學科":      (46,  98, 198),
    "生物科":      (32, 142,  62),
    "生活應用(一)":(198, 112,  22),
    "生活應用(二)":(178,  98,  18),
    "生活應用(三)":(158,  86,  14),
}

WDIR = r"C:\Windows\Fonts"
CDIR = r"C:\Users\user\.claude\plugins\cache\anthropic-agent-skills\document-skills\b0cbd3df1533\skills\canvas-design\canvas-fonts"

def ttf(path, size):
    try:    return ImageFont.truetype(path, size)
    except: return ImageFont.load_default()

f_header = ttf(rf"{WDIR}\kaiu.ttf",  78)
f_huge   = ttf(rf"{WDIR}\kaiu.ttf", 178)
f_sec    = ttf(rf"{WDIR}\kaiu.ttf",  58)
f_award  = ttf(rf"{WDIR}\kaiu.ttf",  52)
f_card   = ttf(rf"{WDIR}\kaiu.ttf",  44)
f_body   = ttf(rf"{WDIR}\kaiu.ttf",  40)
f_sm     = ttf(rf"{WDIR}\kaiu.ttf",  34)
f_xs     = ttf(rf"{WDIR}\kaiu.ttf",  28)

PHOTO_DIR = r"D:\D\onedrive\文件"
awards_data = [
    ("數學科",      "第一名", "你就是那把鎖住個資的鑰匙",  "廖采玲  陳彙喬  黃柏之", "羅婉萍老師",
     rf"{PHOTO_DIR}\數學科第一名與縣長合影.jpg"),
    ("數學科",      "第二名", "平舖直數",                "廖泱泱",              "詩富強老師",
     rf"{PHOTO_DIR}\數學科第二名.jpg"),
    ("生物科",      "第二名", "棘跳蟲對環境壓力的存活力",  "洪靖宇  蔡承恩",      "許議鶴老師",
     rf"{PHOTO_DIR}\生物科第二名.jpg"),
    ("生物科",      "第三名", "斑龜感官行為之綜合研究",    "黃家恩  王嘉裕  陳昱勳","郭靜如老師",
     rf"{PHOTO_DIR}\生物科第三名.jpg"),
    ("生活應用(一)", "第四名", "我的開心智能發電場",       "陳沛青  鍾子鵬  陳意凡","林俊杰老師",
     rf"{PHOTO_DIR}\生活應用一第四名.jpg"),
    ("生活應用(二)", "第三名", "葡萄果乾品質研究",         "籃苡熏  萬奕岑  陳姸汝","許議鶴老師",
     rf"{PHOTO_DIR}\生活與應用科學科(二) 第三名.jpg"),
    ("生活應用(三)", "第四名", "水生生物改善水質研究",     "許呈睿  黃品語  洪啟鈞","許議鶴老師",
     rf"{PHOTO_DIR}\生活與應用科學科(三) 第四名.jpg"),
]

# ── Helpers ────────────────────────────────────────────────────────
def th(draw, text, font):
    bb = draw.textbbox((0,0), text, font=font)
    return bb[3]-bb[1]

def draw_bunting(draw, y0=0, flag_h=88, flag_w=82):
    draw.rectangle([0, 0, W, y0+flag_h], fill=(244, 240, 232))
    draw.line([(0, y0+4),(W, y0+4)], fill=(128, 96, 64), width=5)
    n = W // flag_w + 2
    for i in range(n):
        c   = BUNTING_COLORS[i % len(BUNTING_COLORS)]
        x0_ = i * flag_w
        pts = [(x0_, y0), (x0_+flag_w, y0), (x0_+flag_w//2, y0+flag_h)]
        draw.polygon(pts, fill=c)
        draw.line([(x0_, y0), (x0_+flag_w//2, y0+flag_h)], fill=WHITE, width=2)
        draw.line([(x0_+flag_w, y0), (x0_+flag_w//2, y0+flag_h)], fill=WHITE, width=2)

def draw_grass(draw, start_y):
    pts = [(0, H)]
    for x in range(0, W+1, 20):
        y = start_y + 22*math.sin(x/210) + 11*math.sin(x/68+0.4)
        pts.append((x, int(y)))
    pts.append((W, H))
    draw.polygon(pts, fill=(62, 172, 78))
    pts2 = [(0, H)]
    for x in range(0, W+1, 20):
        y = start_y + 32 + 15*math.sin(x/175+1) + 8*math.sin(x/52+1.1)
        pts2.append((x, int(y)))
    pts2.append((W, H))
    draw.polygon(pts2, fill=(42, 138, 58))

def load_fit(path, w, h, top_bias=0.5):
    """Center-crop photo. top_bias=0 keeps top, 0.5=center, 1=bottom."""
    photo = Image.open(path).convert("RGB")
    pw, ph = photo.size
    scale = max(w/pw, h/ph)
    nw, nh = int(pw*scale)+1, int(ph*scale)+1
    photo = photo.resize((nw, nh), Image.LANCZOS)
    left = (nw-w)//2
    top  = max(0, int((nh-h)*top_bias))
    return photo.crop((left, top, left+w, top+h))

# ── Banner generator ────────────────────────────────────────────────
def make_banner(with_photos=False):
    img  = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(img)

    # Bunting (top 88px)
    BUNT_H = 88
    draw_bunting(draw, y0=0, flag_h=BUNT_H, flag_w=82)

    # Header gradient (88–210)
    HDR_Y0, HDR_H = BUNT_H, 118
    for i in range(HDR_H):
        f = i/HDR_H
        r = int(DARK_GRN[0]*(1-f) + TEAL[0]*f)
        g = int(DARK_GRN[1]*(1-f) + TEAL[1]*f)
        b = int(DARK_GRN[2]*(1-f) + TEAL[2]*f)
        draw.line([(0, HDR_Y0+i),(W, HDR_Y0+i)], fill=(r,g,b))
    draw.line([(0, HDR_Y0+HDR_H),(W, HDR_Y0+HDR_H)], fill=GOLD, width=4)

    hdr = "雲林縣私立正心高級中學  ·  科展榮耀榜"
    bb  = draw.textbbox((0,0), hdr, font=f_header)
    hx  = (W-(bb[2]-bb[0]))//2
    hy  = HDR_Y0 + (HDR_H-(bb[3]-bb[1]))//2 - bb[1]
    draw.text((hx+3, hy+3), hdr, font=f_header, fill=(0,50,30))
    draw.text((hx,   hy),   hdr, font=f_header, fill=WHITE)

    # Content area
    CONT_Y = HDR_Y0 + HDR_H + 14
    CONT_B = H - 70
    AVAIL  = CONT_B - CONT_Y

    draw_grass(draw, CONT_B)

    # Vertical divider
    DIV_X = 1660
    draw.line([(DIV_X, CONT_Y+20),(DIV_X, CONT_B-20)], fill=TEAL, width=3)
    for dy in [CONT_Y+50, H//2, CONT_B-50]:
        draw.ellipse([DIV_X-7, dy-7, DIV_X+7, dy+7], fill=GOLD)

    # ══ LEFT — 高中科展 ═══════════════════════════════════════════════
    LEFT_CX = DIV_X // 2

    huge_h_  = th(draw,"狂賀！",    f_huge)
    sub_h_   = th(draw,"第 66 屆第四區分區高中科展", f_sec)
    at_h_    = th(draw,"榮獲「環境學科」優等", f_award)
    band_h_  = at_h_ + 24
    stud_h_  = th(draw,"高一心班 宋彥霖", f_body)
    teach_h_ = th(draw,"指導老師：張祐誠老師", f_body)
    L_PHOTO_PATH = r"C:\Users\user\Downloads\獲獎片語評審合照.jpg"
    L_PHOTO_W    = DIV_X - 120      # 留左右各 60px 邊距
    L_PHOTO_H    = 520
    LBH = huge_h_+20+sub_h_+22+band_h_+22+stud_h_+16+teach_h_+24+L_PHOTO_H

    y = CONT_Y + (AVAIL-LBH)//2

    draw.text((LEFT_CX+4, y+4), "狂賀！", font=f_huge, fill=(168,48,28), anchor='mt')
    draw.text((LEFT_CX,   y),   "狂賀！", font=f_huge, fill=CORAL,       anchor='mt')
    y += huge_h_ + 20

    draw.text((LEFT_CX,y),"第 66 屆第四區分區高中科展",font=f_sec,fill=NAVY,anchor='mt')
    y += sub_h_ + 22

    at    = "榮獲「環境學科」優等"
    bb_at = draw.textbbox((0,0),at,font=f_award)
    aw    = bb_at[2]-bb_at[0]
    bx0   = LEFT_CX-aw//2-42
    bx1   = LEFT_CX+aw//2+42
    draw.rectangle([bx0,y,bx1,y+band_h_], fill=TEAL, outline=DARK_GRN, width=2)
    draw.text((LEFT_CX,y+band_h_//2),at,font=f_award,fill=WHITE,anchor='mm')
    y += band_h_ + 22

    draw.text((LEFT_CX,y),"高一心班 宋彥霖・高一意班 鐘宥昕・高一正班 廖柃柃",
              font=f_body,fill=DARK,anchor='mt')
    y += stud_h_ + 16
    draw.text((LEFT_CX,y),"指導老師：張祐誠老師",font=f_body,fill=DARK_GRN,anchor='mt')
    y += teach_h_ + 24

    # 左側照片：獲獎與評審合照
    px0 = (DIV_X - L_PHOTO_W) // 2
    try:
        lph = load_fit(L_PHOTO_PATH, L_PHOTO_W, L_PHOTO_H, top_bias=0.12)
        img.paste(lph, (px0, y))
        draw.rectangle([px0-3, y-3, px0+L_PHOTO_W+3, y+L_PHOTO_H+3],
                       outline=TEAL, width=3)
    except Exception as e:
        draw.rectangle([px0, y, px0+L_PHOTO_W, y+L_PHOTO_H], fill=(200,215,225))
        draw.text((LEFT_CX, y+L_PHOTO_H//2), "照片讀取失敗", font=f_sm, fill=DARK, anchor='mm')

    # ══ RIGHT — 縣科展 ════════════════════════════════════════════════
    RX0 = DIV_X + 20
    RW  = W - RX0 - 30
    RCX = RX0 + RW//2

    sec    = "雲林縣第 66 屆公私立國民中小學科學展覽會"
    bb_sec = draw.textbbox((0,0),sec,font=f_sec)
    sec_w_ = bb_sec[2]-bb_sec[0]
    sec_h_ = bb_sec[3]-bb_sec[1]
    ph_    = sec_h_ + 16
    draw.rectangle([RCX-sec_w_//2-28, CONT_Y+8, RCX+sec_w_//2+28, CONT_Y+8+ph_], fill=NAVY)
    draw.text((RCX, CONT_Y+8+ph_//2), sec, font=f_sec, fill=WHITE, anchor='mm')

    CARD_GAP = 14
    CARDS_Y0 = CONT_Y + 8 + ph_ + 16
    CARDS_B  = CONT_B - 4
    cards_h  = CARDS_B - CARDS_Y0
    CARD_H   = (cards_h - CARD_GAP) // 2
    card_w   = (RW - 3*CARD_GAP) // 4

    for i,(subj,rank,topic,students,teacher,photo_path) in enumerate(awards_data):
        if i < 4:
            cx0 = RX0 + i*(card_w+CARD_GAP)
            cy0 = CARDS_Y0
        else:
            col = i-4
            r2w = 3*card_w + 2*CARD_GAP
            r2x = RX0 + (RW-r2w)//2
            cx0 = r2x + col*(card_w+CARD_GAP)
            cy0 = CARDS_Y0 + CARD_H + CARD_GAP

        cx1  = cx0 + card_w
        cy1  = cy0 + CARD_H
        ccx  = (cx0+cx1)//2
        sc   = SUBJ_COLORS.get(subj,(80,80,80))

        # Card background
        draw.rectangle([cx0,cy0,cx1,cy1], fill=WHITE, outline=(205,205,205), width=1)

        if with_photos:
            TEXT_STRIP = 158
            PHOTO_H    = CARD_H - TEXT_STRIP
            try:
                ph_img = load_fit(photo_path, card_w-2, PHOTO_H, top_bias=0.2)
                img.paste(ph_img,(cx0+1, cy0))
            except Exception as e:
                draw.rectangle([cx0+1,cy0,cx1-1,cy0+PHOTO_H], fill=(190,205,220))
                draw.text((ccx,cy0+PHOTO_H//2), "照片讀取失敗", font=f_xs, fill=DARK, anchor='mm')

            # Coloured strip at bottom
            strip_y = cy0 + PHOTO_H
            draw.rectangle([cx0, strip_y, cx1, cy1], fill=sc)
            # 三行各自定位，確保不重疊
            draw.text((ccx, strip_y+44), f"{subj}  {rank}", font=f_card, fill=WHITE, anchor='mm')
            draw.text((ccx, strip_y+86), students, font=f_sm, fill=CREAM, anchor='mt')
            draw.text((ccx, cy1-14), teacher, font=f_xs, fill=(255,255,200), anchor='mb')
        else:
            # Coloured top band
            BAND_H = 72
            draw.rectangle([cx0,cy0,cx1,cy0+BAND_H], fill=sc)
            draw.text((ccx,cy0+BAND_H//2), f"{subj}  {rank}", font=f_card, fill=WHITE, anchor='mm')

            # Content block centred in remaining space
            # G1/G2/G3 match exact gaps used in draw calls below
            G1, G2, G3 = 16, 14, 16
            topic_h_  = th(draw,topic,   f_sm)
            stud_h_c  = th(draw,students, f_sm)
            teach_h_c = th(draw,teacher,  f_sm)
            CONTENT   = topic_h_ + G1 + stud_h_c + G2 + 1 + G3 + teach_h_c
            by0 = cy0 + BAND_H + max(20, (CARD_H - BAND_H - CONTENT)//2)

            draw.text((ccx, by0),                topic,    font=f_sm, fill=(110,110,110), anchor='mt')
            draw.text((ccx, by0+topic_h_+G1),    students, font=f_sm, fill=DARK,         anchor='mt')
            line_y = by0+topic_h_+G1+stud_h_c+G2
            draw.line([(cx0+20,line_y),(cx1-20,line_y)], fill=(210,210,210), width=1)
            draw.text((ccx, line_y+1+G3), teacher, font=f_sm, fill=sc, anchor='mt')

    # Save PNG + PPTX
    suffix   = "照片版" if with_photos else "活潑版"
    out_png  = rf"C:\Users\user\allen\第66屆科展得獎橫幅_{suffix}.png"
    out_pptx = rf"C:\Users\user\allen\第66屆科展得獎橫幅_{suffix}.pptx"
    img.save(out_png, "PNG")

    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width  = Emu(40684450)
    prs.slide_height = Emu(12780963)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(out_png, Emu(0), Emu(0), Emu(40684450), Emu(12780963))
    prs.save(out_pptx)
    print(f"[{suffix}] PNG  → {out_png}")
    print(f"[{suffix}] PPTX → {out_pptx}")

make_banner(with_photos=False)
make_banner(with_photos=True)
print("Done.")
